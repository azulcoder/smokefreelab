"""Render reports/smokefreelab_deck.pptx from reports/deck_outline.md content.

Standalone python-pptx renderer — no pandoc dependency. Produces a 10-slide
16:9 deck with consistent McKinsey-style title/body layout. Speaker notes
from the outline attach to each slide.

Usage:
    python3 scripts/render_deck.py
"""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

REPO_ROOT = Path(__file__).resolve().parents[1]
OUTPUT = REPO_ROOT / "reports" / "smokefreelab_deck.pptx"

INK = RGBColor(0x11, 0x11, 0x11)
SUB = RGBColor(0x55, 0x55, 0x55)
ACCENT = RGBColor(0x8B, 0x00, 0x00)
BAND = RGBColor(0xF3, 0xF1, 0xEE)
RULE = RGBColor(0xCC, 0xCC, 0xCC)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


def add_accent_bar(slide: object) -> None:
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(0.95), Inches(1.2), Inches(0.06)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = ACCENT
    shape.line.fill.background()


def add_footer(slide: object, page: int) -> None:
    tb = slide.shapes.add_textbox(
        Inches(0.5), Inches(7.05), Inches(12.3), Inches(0.3)
    )
    tf = tb.text_frame
    tf.margin_top = Pt(0)
    tf.margin_bottom = Pt(0)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    r = p.add_run()
    r.text = (
        f"SmokeFreeLab  ·  github.com/azulcoder/smokefreelab  ·  "
        f"Ahmad Zulfan (Az)  ·  {page} / 10"
    )
    r.font.size = Pt(9)
    r.font.color.rgb = SUB
    r.font.name = "Helvetica"


def add_title(slide: object, title: str, eyebrow: str | None = None) -> None:
    if eyebrow:
        eb = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.45), Inches(12.3), Inches(0.35)
        )
        ef = eb.text_frame
        p = ef.paragraphs[0]
        r = p.add_run()
        r.text = eyebrow.upper()
        r.font.size = Pt(11)
        r.font.bold = True
        r.font.color.rgb = ACCENT
        r.font.name = "Helvetica"

    tb = slide.shapes.add_textbox(
        Inches(0.5), Inches(1.1), Inches(12.3), Inches(0.8)
    )
    tf = tb.text_frame
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = title
    r.font.size = Pt(28)
    r.font.bold = True
    r.font.color.rgb = INK
    r.font.name = "Helvetica"


def add_bullets(
    slide: object,
    items: list[str],
    top: float = 2.1,
    left: float = 0.5,
    width: float = 12.3,
    height: float = 4.2,
    size: int = 16,
) -> None:
    tb = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    tf = tb.text_frame
    tf.word_wrap = True
    for idx, item in enumerate(items):
        para = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        para.alignment = PP_ALIGN.LEFT
        para.space_after = Pt(8)
        r = para.add_run()
        r.text = "•  " + item
        r.font.size = Pt(size)
        r.font.color.rgb = INK
        r.font.name = "Helvetica"


def add_notes(slide: object, text: str) -> None:
    notes = slide.notes_slide.notes_text_frame
    notes.text = text


def slide_1_cover(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank

    band = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), SLIDE_W, Inches(2.2)
    )
    band.fill.solid()
    band.fill.fore_color.rgb = BAND
    band.line.fill.background()

    accent = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(2.2), SLIDE_W, Inches(0.08)
    )
    accent.fill.solid()
    accent.fill.fore_color.rgb = ACCENT
    accent.line.fill.background()

    tb = slide.shapes.add_textbox(
        Inches(0.8), Inches(2.6), Inches(12), Inches(1.4)
    )
    tf = tb.text_frame
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = "SmokeFreeLab"
    r.font.size = Pt(60)
    r.font.bold = True
    r.font.color.rgb = INK
    r.font.name = "Helvetica"

    tb2 = slide.shapes.add_textbox(
        Inches(0.8), Inches(4.0), Inches(12), Inches(0.6)
    )
    p2 = tb2.text_frame.paragraphs[0]
    r2 = p2.add_run()
    r2.text = "Commercial analytics for a smoke-free product funnel"
    r2.font.size = Pt(22)
    r2.font.color.rgb = SUB
    r2.font.name = "Helvetica"

    tb3 = slide.shapes.add_textbox(
        Inches(0.8), Inches(5.2), Inches(12), Inches(1.6)
    )
    for line, size, color, bold in [
        ("Ahmad Zulfan (Az)", 18, INK, True),
        ("Jakarta · April 2026", 14, SUB, False),
        ("github.com/azulcoder/smokefreelab", 14, ACCENT, True),
    ]:
        p = tb3.text_frame.add_paragraph() if line != "Ahmad Zulfan (Az)" else tb3.text_frame.paragraphs[0]
        r = p.add_run()
        r.text = line
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.color.rgb = color
        r.font.name = "Helvetica"

    add_notes(
        slide,
        "Nine notebooks, one repository, one question — how does a data "
        "scientist defend an IDR budget reallocation. The deck walks through "
        "funnel, experimentation, attribution, MMM, elasticity, and CLV, with "
        "every number framed in rupiah.",
    )


def slide_2_question(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_accent_bar(slide)
    add_title(
        slide,
        "Should we spend the next Rp 50 B / quarter the same way we spent the last one?",
        eyebrow="The question this repo answers",
    )
    add_bullets(
        slide,
        [
            "Is the funnel working end-to-end?",
            "Which channel carries the incremental conversion?",
            "Is the price elastic enough that a cukai hike will drop volume?",
        ],
        top=3.0,
        size=18,
    )
    add_footer(slide, 2)
    add_notes(
        slide,
        "Every slide that follows answers one of these three. If a slide does "
        "not answer one of them, it is cut.",
    )


def slide_3_stack(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_accent_bar(slide)
    add_title(slide, "Nine notebooks, six modules, one quality gate", eyebrow="The stack at a glance")
    add_bullets(
        slide,
        [
            "Nine executable notebooks.  01_ EDA  →  09_ MMM.",
            "Six library modules.  analytics/ · attribution/ · experiment/ · features/ · data/.",
            "One Streamlit app — Experiment Designer, three tabs (Planner / Readout / Peeking).",
            "Quality gate:  ruff + black + mypy --strict + pytest ≥ 60% coverage,  CI green.",
        ],
    )
    add_footer(slide, 3)
    add_notes(
        slide,
        "Portfolio quality bar: typed, tested, linted. The bar is the same as "
        "production; the domain is synthetic.",
    )


def slide_4_funnel(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_accent_bar(slide)
    add_title(slide, "Where the money leaks", eyebrow="Funnel")
    add_bullets(
        slide,
        [
            "354,857 sessions  —  GA4 BigQuery public sample.",
            "20.0% view_item → add_to_cart.  Worst proportional leak ≈ 80% at one stage.",
            "Stage CVR lift economics are asymmetric: +1 pp at stage 2 is worth ≈ 4×  +1 pp at stage 5.",
            "Where we deploy experiments first — the leak with highest rupiah recovery per pp lifted.",
        ],
    )
    add_footer(slide, 4)
    add_notes(
        slide,
        "The 5-minute executive brief: where does the funnel leak, and where do "
        "we deploy experiments for maximum IDR return per experiment.",
    )


def slide_5_experimentation(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_accent_bar(slide)
    add_title(slide, "What to ship when p > 0.05", eyebrow="Experimentation")
    add_bullets(
        slide,
        [
            "SRM gate first.  Reject any experiment that fails SRM at α = 0.01 — lifts cannot be trusted when assignment is broken.",
            "Frequentist CI + Bayesian  P(T > C).  Stakeholders understand  \"P(treatment beats control) = 96%\".  They do not understand  \"p = 0.04\".",
            "Peeking is the real enemy.  Empirical Type-I inflation  5%  →  ≈ 15%  at  10  peeks.  Bayesian posteriors are always-valid.",
        ],
    )
    add_footer(slide, 5)
    add_notes(
        slide,
        "The slide I'd present at the first stakeholder meeting. Re-frame the "
        "conversation from 'is the p-value below 5%' to 'how confident are we "
        "and at what cost if we're wrong'.",
    )


def slide_6_attribution(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_accent_bar(slide)
    add_title(slide, "Where did the conversion come from?", eyebrow="Attribution")
    add_bullets(
        slide,
        [
            "Six methods compared head-to-head:  last-click · first-click · linear · time-decay · Markov · exact Shapley.",
            "Last-click over-credits bottom-funnel channels by ≈ 30–40% vs Shapley.",
            "Compensation schemes pegged to last-click reward conversion, not contribution — brand-building channels are systematically underpaid.",
            "Shapley values are the only credit-allocation scheme satisfying efficiency, symmetry, and null-player simultaneously.",
        ],
    )
    add_footer(slide, 6)
    add_notes(
        slide,
        "Attribution is not an ML problem; it's a coalitional game-theory problem. "
        "Shapley values are the only credit-allocation scheme that satisfies "
        "efficiency, symmetry, and null-player axioms simultaneously.",
    )


def slide_7_mmm(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_accent_bar(slide)
    add_title(slide, "How do we reallocate the next Rp 50 B?", eyebrow="MMM — the signature deliverable")
    add_bullets(
        slide,
        [
            "Current allocation:  TV  32%  ·  digital  12%  ·  trade  52%  ·  other  4%.",
            "Posterior-mean marginal ROI:  TV ≈ 2.4×  ·  digital ≈ 1.9×  ·  trade ≈ 2.7×.",
            "Recommendation:  shift  Rp 5 B  TV → trade this quarter.  Expected gain ≈  Rp 1.4 B  incremental revenue at constant spend.",
            "Risk flag:  trade is near its Hill-saturation elbow.  Beyond  Rp 3 B / week  marginal return drops sharply — monitor weekly.",
        ],
    )
    add_footer(slide, 7)
    add_notes(
        slide,
        "Custom adstock and Hill saturation in PyMC — no black-box package. "
        "Every prior is defensible in this room. The HDI is on the chart, so "
        "readers see the uncertainty, not just the point estimate.",
    )


def slide_8_elasticity(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_accent_bar(slide)
    add_title(slide, "What does the next cukai hike do to volume?", eyebrow="Price elasticity")
    add_bullets(
        slide,
        [
            "Scenario:  Rp 3,000  →  Rp 3,800 per stick  (+26.7%).  Category-wide volume  ≈  −14%  under hierarchical pooling.",
            "Elastic-SKU bucket alone drops  ≈  −33%.  Partial pooling pulls extreme SKU estimates toward the category baseline.",
            "Why hierarchical, not OLS:  thin SKU panels under-cover in OLS.  Pooling stabilises estimates without hiding heterogeneity.",
            "Use both numbers:  OLS as upper bound, hierarchical as central estimate — that is the boundary range CFOs plan against.",
        ],
    )
    add_footer(slide, 8)
    add_notes(
        slide,
        "Elasticity is a core number the finance team asks for every time "
        "cukai moves. Hierarchical partial pooling is the defensible default.",
    )


def slide_9_clv(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_accent_bar(slide)
    add_title(slide, "Who are the top-decile customers", eyebrow="CLV + RFM")
    add_bullets(
        slide,
        [
            "Top 10% of customers by predicted 12-month CLV capture  48%  of aggregate CLV.  Gini ≈ 0.62.",
            "Eleven RFM segments map 1-to-1 to retention tactics:  Champions → VIP · At Risk → win-back · Hibernating → reactivation · Lost → stop spending.",
            "Retention budget  Rp 2 B / quarter  at  3:1 ROAS cap  →  ≈ 73%  of spend to the top two segments.",
            "Which is also where  ≈ 80%  of the returnable CLV sits.  Flat retention programs waste  ≈ 85%  of spend.",
        ],
    )
    add_footer(slide, 9)
    add_notes(
        slide,
        "This slide justifies cutting the long tail and concentrating on "
        "Champions and At-Risk segments.",
    )


def slide_10_impact(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_accent_bar(slide)
    add_title(slide, "Rp 27 B / year incremental CLV on a 250K-registrant funnel", eyebrow="Impact and close")
    add_bullets(
        slide,
        [
            "First 90 days in a commercial analytics role:  wire POS + GA4 + commerce data into the modules above  (weeks 1–3).",
            "Launch three A/B experiments through the Experiment Designer workflow  (weeks 4–8).",
            "Ship first rupiah-denominated MMM readout with HDI to Brand Directors  (weeks 9–12).",
            "Repo:  github.com/azulcoder/smokefreelab.",
        ],
        size=15,
    )
    add_footer(slide, 10)
    add_notes(
        slide,
        "The Rp 27 B number is conservative. The full sensitivity matrix — "
        "1 pp to 3 pp lift × Rp 225K to Rp 675K LTV — spans Rp 6.75 B to "
        "Rp 60 B. The point is not the point estimate; the point is I can "
        "tell you where the number falls on that grid given an input I trust.",
    )


def main() -> None:
    OUTPUT.parent.mkdir(parents=True, exist_ok=True)
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    slide_1_cover(prs)
    slide_2_question(prs)
    slide_3_stack(prs)
    slide_4_funnel(prs)
    slide_5_experimentation(prs)
    slide_6_attribution(prs)
    slide_7_mmm(prs)
    slide_8_elasticity(prs)
    slide_9_clv(prs)
    slide_10_impact(prs)

    prs.core_properties.title = "SmokeFreeLab"
    prs.core_properties.author = "Ahmad Zulfan (Az)"
    prs.core_properties.subject = "Commercial analytics playbook for a smoke-free product funnel"

    prs.save(str(OUTPUT))
    print(f"Wrote {OUTPUT}  ({len(prs.slides)} slides)")


if __name__ == "__main__":
    main()
